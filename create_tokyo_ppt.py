from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

# Tokyo Night Colors
BG_COLOR = RGBColor(26, 27, 38)      # #1a1b26 (Deep Blue)
FG_COLOR = RGBColor(192, 202, 245)   # #c0caf5 (Pale Blue Text)
PRIMARY = RGBColor(122, 162, 247)    # #7aa2f7 (Blue)
SECONDARY = RGBColor(187, 154, 247)  # #bb9af7 (Purple)
SUCCESS = RGBColor(158, 206, 106)    # #9ece6a (Green)
CARD_BG = RGBColor(36, 40, 59)       # #24283b (Card)

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_card(slide, left, top, width, height, title_text, items):
    # Card Background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(65, 72, 104) # Lighter border
    
    # Text Frame
    tf = shape.text_frame
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.2)
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.color.rgb = FG_COLOR
        p.space_before = Pt(6)

def create_tokyo_ppt():
    prs = Presentation()
    
    # 1. Cover Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    set_slide_background(slide, BG_COLOR)
    
    # Decorative elements
    shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(-1), Inches(4), Inches(4), Inches(4))
    shape.rotation = 45
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()

    title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    p = title.text_frame.paragraphs[0]
    p.text = "Python Email Checker"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
    p = subtitle.text_frame.paragraphs[0]
    p.text = "이메일 업무 자동화 시스템 구조 분석\n2026-01-01"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(24)
    p.font.color.rgb = SECONDARY

    # 2. System Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    
    # Header
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "1. 시스템 개요"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # Cards
    card_width = Inches(3)
    card_height = Inches(3.5)
    margin = Inches(0.3)
    start_left = Inches(0.5)
    top = Inches(2)
    
    add_card(slide, start_left, top, card_width, card_height, "이메일 모니터링", 
             ["IMAP 프로토콜 활용", "60초 주기 자동 감지", "읽지 않은 메일 수집"])
             
    add_card(slide, start_left + card_width + margin, top, card_width, card_height, "지능형 필터링", 
             ["설정된 날짜 이후만 감지", "발신자 목록 대조", "제목 키워드 검색"])
             
    add_card(slide, start_left + (card_width + margin)*2, top, card_width, card_height, "실시간 알림", 
             ["텔레그램 Bot API 연동", "즉각적인 푸시 알림", "처리 결과 로깅"])

    # 3. Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "2. 아키텍처 다이어그램"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    img_path = 'tokyo_diagram.png'
    if os.path.exists(img_path):
        pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0))
        
        # Resize to fit (Target height 5.5 inches)
        target_height = Inches(5.5)
        ratio = target_height / pic.height
        pic.height = int(pic.height * ratio)
        pic.width = int(pic.width * ratio)
        
        # Center
        slide_width = prs.slide_width
        pic.left = int((slide_width - pic.width) / 2)
        pic.top = Inches(1.5)
        
        # Border
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, pic.left, pic.top, pic.width, pic.height)
        line.fill.background()
        line.line.color.rgb = SECONDARY
        line.line.width = Pt(2)

    # 4. Core Components
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "3. 핵심 컴포넌트"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    add_card(slide, Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5), "메인 컨트롤러 (Main)", 
             ["진입점(Entry Point): main.py", "스케줄러 작업 관리", "테스트 모드 (--test) 지원"])
             
    add_card(slide, Inches(5.2), Inches(1.5), Inches(4.5), Inches(2.5), "체크 & 필터 (Checker)", 
             ["서버 측 필터: 날짜 조건", "클라이언트 필터: 발신자/키워드", "중복 처리 방지 (UID)"])

    add_card(slide, Inches(0.5), Inches(4.2), Inches(4.5), Inches(2.5), "알림 서비스 (Notifier)", 
             ["Requests 라이브러리 사용", "HTTP POST 요청 전송", "예외 처리 및 로깅"])

    add_card(slide, Inches(5.2), Inches(4.2), Inches(4.5), Inches(2.5), "환경 설정 (Config)", 
             [".env 파일 로드 (보안)", "설정값 전처리 (List 변환)", "중앙 집중식 관리"])

    # 5. Data Flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_COLOR)
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "4. 데이터 흐름 (Data Flow)"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    steps = [
        "1. 설정 로드: .env 파일 및 환경변수 초기화",
        "2. 스케줄링 트리거: 60초마다 Job 실행",
        "3. 데이터 수집: IMAP 서버 접속 및 안 읽은 메일 확인",
        "4. 필터링 적용: 날짜 -> 발신자 -> 키워드 검사",
        "5. 알림 발송: 조건 일치 시 텔레그램 메시지 전송"
    ]
    
    top = Inches(1.5)
    for i, step in enumerate(steps):
        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(1), top, Inches(8), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = SUCCESS if i == 4 else SECONDARY
        shape.line.width = Pt(1.5)
        
        p = shape.text_frame.paragraphs[0]
        p.text = step
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(18)
        p.font.color.rgb = FG_COLOR
        p.alignment = PP_ALIGN.LEFT
        shape.text_frame.margin_left = Inches(0.5)
        
        top += Inches(1.0)

    prs.save('tokyo_analysis_report.pptx')
    print("Tokyo Night PPT generated: tokyo_analysis_report.pptx")

if __name__ == "__main__":
    create_tokyo_ppt()
