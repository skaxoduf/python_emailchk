from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, font_size=40, is_bold=True, is_white=True):
    title = slide.shapes.title
    title.text = text
    p = title.text_frame.paragraphs[0]
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(font_size)
    p.font.bold = is_bold
    
    if is_white:
        p.font.color.rgb = RGBColor(255, 255, 255)
    else:
        p.font.color.rgb = RGBColor(0, 191, 255) # Deep Sky Blue

def add_card(slide, left, top, width, height, title_text, items):
    # Card Background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(45, 45, 48) # Dark Gray
    shape.line.color.rgb = RGBColor(60, 60, 60)
    
    # Text Frame
    tf = shape.text_frame
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.2)
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 191, 255) # Cyan Accent
    
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.space_before = Pt(6)

def create_modern_ppt():
    prs = Presentation()
    
    # Colors
    BG_COLOR = RGBColor(30, 30, 30) # Almost Black
    ACCENT_COLOR = RGBColor(0, 191, 255) # Deep Sky Blue
    WHITE = RGBColor(255, 255, 255)

    # 1. Cover Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    set_slide_background(slide, BG_COLOR)
    
    title = slide.shapes.title
    title.text = "Python Email Checker"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.name = "Malgun Gothic"
    title.text_frame.paragraphs[0].font.size = Pt(54)

    subtitle = slide.placeholders[1]
    subtitle.text = "시스템 아키텍처 및 코드 분석 보고서\n2026-01-01"
    for p in subtitle.text_frame.paragraphs:
        p.font.color.rgb = ACCENT_COLOR
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(24)

    # 2. Overview Slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BG_COLOR)
    
    # Header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.1), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "1. System Overview"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Cards
    card_width = Inches(3)
    card_height = Inches(3.5)
    margin = Inches(0.3)
    start_left = Inches(0.5)
    top = Inches(2)
    
    add_card(slide, start_left, top, card_width, card_height, "Email Monitoring", 
             ["IMAP 프로토콜 활용", "주기적(60s) 실행", "읽지 않은 메일 감지"])
             
    add_card(slide, start_left + card_width + margin, top, card_width, card_height, "Smart Filtering", 
             ["날짜 기반 1차 필터", "발신자 목록 확인", "제목 키워드 매칭"])
             
    add_card(slide, start_left + (card_width + margin)*2, top, card_width, card_height, "Instant Alert", 
             ["텔레그램 Bot API", "실시간 푸시 알림", "전송 결과 로깅"])

    # 3. Architecture Slide
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BG_COLOR)
    
    # Header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.1), Inches(0.6))
    shape.fill.solid() # Fixed: fill.solid() needs to be called on fill object
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "2. Architecture Diagram"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    img_path = 'architecture_diagram.png'
    if os.path.exists(img_path):
        # Center the image
        # Slide width is 10 inches, height is 7.5 inches
        # Image we want roughly 8x6 or proportional
        
        pic = slide.shapes.add_picture(img_path, Inches(0), Inches(0))
        
        # Calculate centering
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Resize image to fit keeping aspect ratio
        # Target height: 5.5 inches
        target_height = Inches(5.5)
        ratio = target_height / pic.height
        pic.height = int(pic.height * ratio)
        pic.width = int(pic.width * ratio)
        
        # Center position
        pic.left = int((slide_width - pic.width) / 2)
        pic.top = Inches(1.5) # Below title
        
        # Add a border effect (manual line)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, pic.left, pic.top, pic.width, pic.height)
        line.fill.background() # No fill
        line.line.color.rgb = ACCENT_COLOR
        line.line.width = Pt(2)
    else:
        tb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(5), Inches(1))
        tb.text = "Image Not Found!"
        tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)

    # 4. Core Components
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BG_COLOR)
    
    # Header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.1), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "3. Core Components"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Logic
    add_card(slide, Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5), "Main Controller", 
             ["Entry Point: main.py", "스케줄러 관리 및 예외 처리", "테스트 모드 지원"])
             
    add_card(slide, Inches(5.2), Inches(1.5), Inches(4.5), Inches(2.5), "Email Checker Filter", 
             ["Server-side: Date filtering", "Client-side: Sender & Keyword", "Dual-Layer Filtering Logic"])

    add_card(slide, Inches(0.5), Inches(4.2), Inches(4.5), Inches(2.5), "Notifier Service", 
             ["Requests Library 활용", "HTTP POST to Telegram", "Time-out 및 에러 핸들링"])

    add_card(slide, Inches(5.2), Inches(4.2), Inches(4.5), Inches(2.5), "Configuration", 
             ["Environment Variables (.env)", "Secure Credential Management", "Data Parsing & Validation"])

    # 5. Data Flow (Process)
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BG_COLOR)
    
    # Header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.1), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "4. Data Flow Process"
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    steps = [
        "1. Start & Load Config (.env)",
        "2. Scheduled Trigger (Every 60s)",
        "3. Fetch Unseen Emails (IMAP)",
        "4. Apply Custom Filters (Date/Sender/Key)",
        "5. Push Notification (Telegram API)"
    ]
    
    top = Inches(1.5)
    for i, step in enumerate(steps):
        # Arrow Shape
        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(1), top, Inches(8), Inches(0.8))
        shape.fill.solid()
        if i % 2 == 0:
            shape.fill.fore_color.rgb = RGBColor(50, 50, 55)
        else:
            shape.fill.fore_color.rgb = RGBColor(60, 60, 65)
        shape.line.color.rgb = ACCENT_COLOR
        shape.line.width = Pt(1.5)
        
        p = shape.text_frame.paragraphs[0]
        p.text = step
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        shape.text_frame.margin_left = Inches(0.5)
        
        top += Inches(1.0)

    # Save
    output_filename = 'modern_analysis_report.pptx'
    prs.save(output_filename)
    print(f"Modern PPTX created: {output_filename}")

if __name__ == "__main__":
    create_modern_ppt()
